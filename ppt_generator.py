"""
PPT_GENERATOR.PY — Landsat-v1
Genera PPT con timelapses RGB + SWIR a partir de la plantilla Copernicus.
Busca texto "color verdadero" -> RGB  y  "falso color" -> SWIR en la plantilla.
"""
import os, glob, shutil, zipfile, tempfile, re
from lxml import etree
from datetime import datetime
from pptx import Presentation
from PIL import Image
from gif_cache import existe_en_cache
from config_landsat import VOLCANES

VOLCANES_ACTIVOS = list(VOLCANES.keys())
PLANTILLA_PATH   = "docs/plantillas/Cambios_morfologicos.pptx"
OUTPUT_DIR       = "docs/landsat"

MESES_ES = {1:'enero',2:'febrero',3:'marzo',4:'abril',5:'mayo',6:'junio',
            7:'julio',8:'agosto',9:'septiembre',10:'octubre',11:'noviembre',12:'diciembre'}

def formatear_fecha_espanol(fecha_str):
    try:
        dt = datetime.strptime(fecha_str,'%Y-%m-%d')
        return f"{dt.strftime('%d')} {MESES_ES[dt.month]}"
    except Exception:
        return fecha_str

def comprimir_gif(src, dst, max_mb=1.0):
    size_mb = os.path.getsize(src)/1048576
    if size_mb <= max_mb:
        shutil.copy2(src, dst)
        print(f"      OK ({size_mb:.2f} MB)")
        return dst
    print(f"      Comprimiendo ({size_mb:.2f} MB)...")
    try:
        img = Image.open(src)
        frames = []
        try:
            while True:
                frames.append(img.copy().convert('P', palette=Image.ADAPTIVE, colors=64))
                img.seek(img.tell()+1)
        except EOFError:
            pass
        if size_mb > 2.0:
            ns = (int(frames[0].width*0.75), int(frames[0].height*0.75))
            frames = [f.resize(ns, Image.Resampling.LANCZOS) for f in frames]
        frames[0].save(dst,save_all=True,append_images=frames[1:],
                       optimize=True,duration=img.info.get('duration',1000),loop=0)
        print(f"      {size_mb:.2f}MB -> {os.path.getsize(dst)/1048576:.2f}MB")
    except Exception as e:
        print(f"      Error comprimiendo: {e}")
        shutil.copy2(src, dst)
    return dst

def _set_run_text(shape, new_text):
    """Reemplaza el texto del primer parrafo de un shape preservando formato."""
    p = shape.text_frame.paragraphs[0]
    fmt = None
    if p.runs:
        r0 = p.runs[0]
        fmt = {'name':r0.font.name,'size':r0.font.size,
               'bold':r0.font.bold,'italic':r0.font.italic}
    p.clear()
    run = p.add_run()
    run.text = new_text
    if fmt:
        if fmt['name']:   run.font.name   = fmt['name']
        if fmt['size']:   run.font.size   = fmt['size']
        if fmt['bold']  is not None: run.font.bold   = fmt['bold']
        if fmt['italic'] is not None: run.font.italic = fmt['italic']

def generar_ppt(volcan):
    print(f"\n--- {volcan}")
    carpeta = f"docs/landsat/{volcan}/timelapses_ppt"
    if not os.path.exists(carpeta):
        print(f"    No existe: {carpeta}"); return None

    gifs_rgb  = sorted(glob.glob(f"{carpeta}/{volcan}_RGB_*.gif"))
    gifs_swir = sorted(glob.glob(f"{carpeta}/{volcan}_SWIR_*.gif"))
    print(f"    GIFs: RGB={len(gifs_rgb)}, SWIR={len(gifs_swir)}")
    if not gifs_rgb or not gifs_swir:
        print("    GIFs incompletos (necesita RGB + SWIR)"); return None

    gif_rgb  = gifs_rgb[-1]
    gif_swir = gifs_swir[-1]
    partes   = os.path.basename(gif_rgb).replace('.gif','').split('_')
    if len(partes) < 4: print("    Nombre de GIF inesperado"); return None

    fecha_inicio, fecha_fin = partes[-2], partes[-1]
    print(f"    {fecha_inicio} -> {fecha_fin}")

    temp_rgb  = f"/tmp/{volcan}_RGB.gif"
    temp_swir = f"/tmp/{volcan}_SWIR.gif"
    gif_rgb_final  = comprimir_gif(gif_rgb,  temp_rgb)
    gif_swir_final = comprimir_gif(gif_swir, temp_swir)

    if not os.path.exists(PLANTILLA_PATH):
        print(f"    Plantilla no encontrada: {PLANTILLA_PATH}"); return None

    prs   = Presentation(PLANTILLA_PATH)
    slide = prs.slides[0]

    f_ini_es = formatear_fecha_espanol(fecha_inicio)
    f_fin_es = formatear_fecha_espanol(fecha_fin)
    ano      = fecha_fin.split('-')[0]
    texto_rgb  = f"Imagenes Landsat 8/9 color natural, Time Lapse {f_ini_es} a {f_fin_es} {ano}"
    texto_swir = f"Imagenes Landsat 8/9 SWIR (B7-B6-B4), Time Lapse {f_ini_es} a {f_fin_es} {ano}"

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"): continue
        texto = shape.text.strip().lower()
        if "color verdadero" in texto and "time lapse" in texto:
            _set_run_text(shape, texto_rgb)
        elif "falso color" in texto and "time lapse" in texto:
            _set_run_text(shape, texto_swir)
        elif "volcan" in texto or "volcano" in texto:
            for va in VOLCANES_ACTIVOS:
                if va.lower() in texto:
                    for p in shape.text_frame.paragraphs:
                        tp = p.text
                        if va.lower() in tp.lower():
                            patron = re.compile(re.escape(va), re.IGNORECASE)
                            nuevo  = patron.sub(volcan, tp)
                            try:
                                dt = datetime.strptime(fecha_fin,'%Y-%m-%d')
                                nuevo = re.sub(r'(mes de )\w+(?: \d{4})?',
                                               f"mes de {MESES_ES[dt.month]} {dt.year}",
                                               nuevo, flags=re.IGNORECASE)
                            except Exception:
                                pass
                            if nuevo != tp:
                                fmt = None
                                if p.runs:
                                    r0 = p.runs[0]
                                    fmt = {'name':r0.font.name,'size':r0.font.size,
                                           'bold':r0.font.bold,'italic':r0.font.italic}
                                p.clear(); run = p.add_run(); run.text = nuevo
                                if fmt:
                                    if fmt['name']:   run.font.name   = fmt['name']
                                    if fmt['size']:   run.font.size   = fmt['size']
                                    if fmt['bold']  is not None: run.font.bold   = fmt['bold']
                                    if fmt['italic'] is not None: run.font.italic = fmt['italic']

    # Emparejar imagenes con subtitulos por proximidad horizontal
    subtitle_rgb = subtitle_swir = None
    for s in slide.shapes:
        if not s.has_text_frame: continue
        t = s.text_frame.text.lower()
        if "color verdadero" in t and "time lapse" in t:   subtitle_rgb  = s
        elif "falso color"   in t and "time lapse" in t:   subtitle_swir = s
    shapes_img = [s for s in slide.shapes if s.shape_type == 13]

    def cx_dist(pic, sub):
        return abs((pic.left+pic.width/2)-(sub.left+sub.width/2))

    if len(shapes_img) >= 2 and subtitle_rgb and subtitle_swir:
        pic_rgb  = min(shapes_img, key=lambda p: cx_dist(p, subtitle_rgb))
        pic_swir = next(p for p in shapes_img if p is not pic_rgb)
        for pic, gif_final, label in [(pic_rgb,gif_rgb_final,'RGB'),(pic_swir,gif_swir_final,'SWIR')]:
            pos = (pic.left,pic.top,pic.width,pic.height)
            pic.element.getparent().remove(pic.element)
            slide.shapes.add_picture(gif_final,*pos)
            print(f"       {label} OK")
    elif len(shapes_img) >= 2:
        sorted_imgs = sorted(shapes_img, key=lambda s:(s.top,s.left))
        for pic, gif_final, label in zip(sorted_imgs,[gif_rgb_final,gif_swir_final],['RGB','SWIR']):
            pos = (pic.left,pic.top,pic.width,pic.height)
            pic.element.getparent().remove(pic.element)
            slide.shapes.add_picture(gif_final,*pos)
            print(f"       {label} OK (fallback)")

    carpeta_rep = f"docs/landsat/{volcan}/reportes"
    os.makedirs(carpeta_rep, exist_ok=True)
    out = f"{carpeta_rep}/{volcan}_Evaluacion_Mensual_{fecha_fin[:7]}.pptx"
    try:
        prs.save(out)
        size_mb = os.path.getsize(out)/1048576
        print(f"   PPT: {size_mb:.2f} MB -> {out}")
        for t in [temp_rgb, temp_swir]:
            try: os.remove(t)
            except Exception: pass
        return out
    except Exception as e:
        print(f"    Error guardando PPT: {e}"); return None

def generar_ppt_combinado(ppts, fecha_inicio, fecha_fin):
    """Combina PPTs individuales en un archivo via manipulacion ZIP directa."""
    if not ppts: return None
    print(f"\nGENERANDO PPT COMBINADO ({len(ppts)} volcanes)...")
    os.makedirs("docs/reportes", exist_ok=True)
    output  = f"docs/reportes/Evaluacion_Completa_{fecha_inicio}_{fecha_fin}.pptx"
    tmpdir  = tempfile.mkdtemp()
    NS_PRS  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NS_R    = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    NS_CT   = 'http://schemas.openxmlformats.org/package/2006/content-types'
    NS_RELS = 'http://schemas.openxmlformats.org/package/2006/relationships'
    try:
        with zipfile.ZipFile(ppts[0],'r') as z: z.extractall(tmpdir)
        prs_xml  = os.path.join(tmpdir,'ppt','presentation.xml')
        prs_rels = os.path.join(tmpdir,'ppt','_rels','presentation.xml.rels')
        ct_path  = os.path.join(tmpdir,'[Content_Types].xml')
        prs_tree = etree.parse(prs_xml)
        sldIdLst = prs_tree.find(f'{{{NS_PRS}}}sldIdLst')
        next_sn  = len(sldIdLst)+1
        next_sid = max(int(el.get('id')) for el in sldIdLst)+1
        media_dir= os.path.join(tmpdir,'ppt','media'); os.makedirs(media_dir,exist_ok=True)
        next_mn  = len(os.listdir(media_dir))+1
        for ppt_path in ppts[1:]:
            try:
                with zipfile.ZipFile(ppt_path,'r') as z:
                    nms = z.namelist()
                    slides_z = sorted([n for n in nms if n.startswith('ppt/slides/slide')
                                       and n.endswith('.xml') and '_rels' not in n])
                    if not slides_z: continue
                    sz = slides_z[0]
                    rz = f"ppt/slides/_rels/{os.path.basename(sz)}.rels"
                    sxml = z.read(sz)
                    try: rbytes = z.read(rz)
                    except KeyError: rbytes = None
                    mm = {}
                    for item in nms:
                        if item.startswith('ppt/media/'):
                            old = os.path.basename(item); ext = os.path.splitext(old)[1]
                            new = f"media{next_mn}{ext}"; next_mn += 1; mm[old] = new
                            with z.open(item) as src, open(os.path.join(media_dir,new),'wb') as dst:
                                dst.write(src.read())
                    if rbytes and mm:
                        rt = etree.fromstring(rbytes)
                        for rel in rt.findall(f'{{{NS_RELS}}}Relationship'):
                            t = rel.get('Target',''); on = os.path.basename(t)
                            if on in mm: rel.set('Target',f'../media/{mm[on]}')
                        rbytes = etree.tostring(rt,xml_declaration=True,
                                                encoding='UTF-8',standalone=True)
                sn = f"slide{next_sn}.xml"
                with open(os.path.join(tmpdir,'ppt','slides',sn),'wb') as f: f.write(sxml)
                if rbytes:
                    rd = os.path.join(tmpdir,'ppt','slides','_rels'); os.makedirs(rd,exist_ok=True)
                    with open(os.path.join(rd,f"{sn}.rels"),'wb') as f: f.write(rbytes)
                rId = f"rId_slide{next_sn}"
                el = etree.SubElement(sldIdLst,f'{{{NS_PRS}}}sldId',id=str(next_sid))
                el.set(f'{{{NS_R}}}id',rId)
                rtp = etree.parse(prs_rels)
                etree.SubElement(rtp.getroot(),f'{{{NS_RELS}}}Relationship',Id=rId,
                    Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide',
                    Target=f'slides/{sn}')
                rtp.write(prs_rels,xml_declaration=True,encoding='UTF-8',standalone=True)
                ctt = etree.parse(ct_path)
                etree.SubElement(ctt.getroot(),f'{{{NS_CT}}}Override',
                    PartName=f'/ppt/slides/{sn}',
                    ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
                ctt.write(ct_path,xml_declaration=True,encoding='UTF-8',standalone=True)
                next_sn += 1; next_sid += 1
                print(f"   OK: {os.path.basename(ppt_path)}")
            except Exception as e:
                print(f"   Error {os.path.basename(ppt_path)}: {e}")
        prs_tree.write(prs_xml,xml_declaration=True,encoding='UTF-8',standalone=True)
        if os.path.exists(output): os.remove(output)
        with zipfile.ZipFile(output,'w',zipfile.ZIP_DEFLATED) as zout:
            for root,dirs,files in os.walk(tmpdir):
                for file in files:
                    fp = os.path.join(root,file)
                    zout.write(fp, os.path.relpath(fp,tmpdir))
        size_mb = os.path.getsize(output)/1048576
        print(f"PPT combinado: {output} ({size_mb:.1f} MB)")
        return output
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

def main():
    print("="*60)
    print("PPT GENERATOR — Landsat-v1")
    print("="*60)
    if not os.path.exists(PLANTILLA_PATH):
        print(f"Plantilla no encontrada: {PLANTILLA_PATH}"); return
    volcan_env   = os.getenv('VOLCAN')
    fecha_inicio = os.getenv('FECHA_INICIO')
    fecha_fin    = os.getenv('FECHA_FIN')
    volcanes     = [volcan_env] if volcan_env else VOLCANES_ACTIVOS
    ppts = []
    for v in volcanes:
        try:
            ppt = generar_ppt(v)
            if ppt: ppts.append(ppt)
        except Exception as e:
            print(f"Error {v}: {e}")
    print(f"\n{len(ppts)} PPTs individuales")
    if len(ppts) > 1:
        if not fecha_inicio or not fecha_fin:
            partes = os.path.basename(ppts[0]).replace('.pptx','').split('_')
            fecha_inicio = partes[-2] if len(partes)>=2 else 'inicio'
            fecha_fin    = partes[-1] if len(partes)>=1 else 'fin'
        generar_ppt_combinado(ppts, fecha_inicio, fecha_fin)

if __name__ == "__main__":
    main()
